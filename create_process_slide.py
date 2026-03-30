"""
Generate a BCG/McKinsey-style process flow slide with findings.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --- Constants ---
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Color palette (dark, muted consulting style)
DARK_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
MEDIUM_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
ARROW_GRAY = RGBColor(0xB0, 0xB8, 0xC4)
BORDER_LINE = RGBColor(0xD5, 0xDB, 0xE1)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# --- Background ---
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.08
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=10, font_color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_arrow(slide, left, top, width, height):
    """Add a right-pointing chevron arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ARROW_GRAY
    shape.line.fill.background()
    shape.rotation = 0
    return shape


# ============================================================
# HEADER SECTION
# ============================================================
# Top accent line
accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
accent_line.fill.solid()
accent_line.fill.fore_color.rgb = DARK_NAVY
accent_line.line.fill.background()

# Title
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(10), Inches(0.55),
             "Processkarta: Inköp-till-betalning (P2P)", font_size=22, font_color=DARK_NAVY, bold=True)

# Subtitle
add_text_box(slide, Inches(0.7), Inches(0.72), Inches(10), Inches(0.35),
             "Identifierade findings och deras koppling till processflödet", font_size=12, font_color=MED_GRAY)

# Thin separator line
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.1), Inches(11.93), Pt(1))
sep.fill.solid()
sep.fill.fore_color.rgb = BORDER_LINE
sep.line.fill.background()

# ============================================================
# PROCESS FLOW (7 steps + Uppföljning bar underneath)
# ============================================================
process_steps = [
    "Behov\nuppstår",
    "Upphandling\n/ avtal",
    "Beställning",
    "Mottagning",
    "Faktura",
    "Attest",
    "Betalning",
]

num_steps = len(process_steps)
flow_top = Inches(1.35)
step_w = Inches(1.30)
step_h = Inches(0.72)
arrow_w = Inches(0.32)
arrow_h = Inches(0.28)
total_flow_w = num_steps * step_w + (num_steps - 1) * arrow_w
flow_left_start = (SLIDE_W - total_flow_w) / 2

step_positions = []  # store (left, center, right) for each step

for i, step_text in enumerate(process_steps):
    x = flow_left_start + i * (step_w + arrow_w)

    # Step box
    box = add_rounded_rect(slide, x, flow_top, step_w, step_h, DARK_NAVY, border_color=None)
    box.adjustments[0] = 0.12
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = step_text
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    txBody = tf._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')

    center_x = x + step_w / 2
    step_positions.append((x, center_x, x + step_w))

    # Arrow between steps
    if i < num_steps - 1:
        arrow_x = x + step_w + Emu(2000)
        arrow_y = flow_top + (step_h - arrow_h) / 2
        add_arrow(slide, arrow_x, arrow_y, arrow_w - Emu(4000), arrow_h)

# ============================================================
# "UPPFÖLJNING" — full-width bar beneath the process flow
# ============================================================
uppf_top = flow_top + step_h + Inches(0.18)
uppf_left = step_positions[0][0]
uppf_right = step_positions[-1][2]
uppf_width = uppf_right - uppf_left
uppf_h = Inches(0.38)

uppf_bar = add_rounded_rect(slide, uppf_left, uppf_top, uppf_width, uppf_h,
                             DARK_NAVY, border_color=None)
uppf_bar.adjustments[0] = 0.15
tf = uppf_bar.text_frame
tf.word_wrap = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = "Uppföljning"
p.font.size = Pt(10)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Calibri"
bodyPr = tf._txBody.find(qn('a:bodyPr'))
bodyPr.set('anchor', 'ctr')


# ============================================================
# FINDINGS SECTION (1-6 map to process steps, 7 maps to Uppföljning bar)
# ============================================================
findings = [
    {
        "title": "Ingen styrd väg från behov till köp",
        "steps": (0, 2),  # Behov → Beställning
        "desc": "Det finns en ingång för behov, men inte ett styrt beställningsflöde som leder till rätt leverantör, rätt avtal och rätt pris.",
        "num": "1",
    },
    {
        "title": "Svagt stöd för köp på avtal",
        "steps": (1, 2),  # Upphandling → Beställning
        "desc": "Det är svårt att hitta och använda rätt avtal i vardagen, vilket driver köp utanför avtal.",
        "num": "2",
    },
    {
        "title": "Splittrad avtals- och leverantörsinformation",
        "steps": (1, 6),  # Upphandling → Betalning (full span)
        "desc": "Ingen sammanhållen koppling mellan inköp, avtal och ekonomi försvårar spårbarhet och uppföljning.",
        "num": "3",
    },
    {
        "title": "Manuella och svårt överblickbara arbetsflöden",
        "steps": (2, 6),  # Beställning → Betalning
        "desc": "Excel, mejl, Teams och personberoende arbetssätt försvårar översikt, status och skalbarhet.",
        "num": "4",
    },
    {
        "title": "Otydliga ansvar och varierande arbetssätt",
        "steps": (3, 6),  # Mottagning → Betalning
        "desc": "Roller arbetar olika och ansvar är inte alltid tydligt definierade, vilket försvårar återkoppling och konsekvent hantering.",
        "num": "5",
    },
    {
        "title": "Kontrollen kommer först vid fakturan",
        "steps": (4, 6),  # Faktura → Betalning
        "desc": "Avvikelser och fel upptäcks sent, vilket gör kontrollen reaktiv och försvårar internkontroll och tidig felupptäckt.",
        "num": "6",
    },
]

# Finding 7 is special — it maps to the full-width "Uppföljning" bar
finding_7 = {
    "title": "Splittrad data för budget, prognos och uppföljning",
    "desc": "Utfall, personal, verksamhet och planering finns i flera system och i Excel, vilket försvårar analys, prognoser och scenarier.",
    "num": "7",
}

# Layout settings
findings_start_y = Inches(2.65)
finding_row_h = Inches(0.58)
finding_gap = Inches(0.08)

NUM_MARGIN_LEFT = Inches(0.25)
num_circle_size = Inches(0.32)

for idx, f in enumerate(findings):
    y = findings_start_y + idx * (finding_row_h + finding_gap)
    start_step = f["steps"][0]
    end_step = f["steps"][1]

    # Card spans from left edge of first step to right edge of last step
    card_right = step_positions[end_step][2]
    natural_left = step_positions[start_step][0]
    natural_width = card_right - natural_left
    min_card_width = Inches(5.5)
    min_left = Inches(0.65)
    if natural_width < min_card_width:
        card_left = max(min_left, card_right - min_card_width)
        card_width = card_right - card_left
    else:
        card_width = natural_width
        card_left = natural_left
    card_h = finding_row_h

    # Finding card
    card = add_rounded_rect(slide, card_left, y, card_width, card_h,
                            LIGHT_GRAY, border_color=DARK_NAVY, border_width=Pt(1.2))
    card.adjustments[0] = 0.06

    # Number circle
    num_x = NUM_MARGIN_LEFT
    num_y = y + (card_h - num_circle_size) / 2

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, num_x, num_y, num_circle_size, num_circle_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK_NAVY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = f["num"]
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    bodyPr.set('lIns', '0')
    bodyPr.set('rIns', '0')
    bodyPr.set('tIns', '0')
    bodyPr.set('bIns', '0')

    # Text inside card
    inner_pad = Inches(0.15)
    text_left = card_left + inner_pad
    text_top = y + Inches(0.06)
    text_h = card_h - Inches(0.12)

    if card_width < Inches(4.0):
        # Narrow card: stacked vertically
        txBox = slide.shapes.add_textbox(text_left, text_top, card_width - 2 * inner_pad, text_h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = f["title"]
        p1.font.size = Pt(9)
        p1.font.color.rgb = DARK_NAVY
        p1.font.bold = True
        p1.font.name = "Calibri"
        p1.space_after = Pt(2)
        p2 = tf.add_paragraph()
        p2.text = f["desc"]
        p2.font.size = Pt(7.5)
        p2.font.color.rgb = MED_GRAY
        p2.font.name = "Calibri"
    else:
        # Wide card: title left, description right
        title_w = min(Inches(3.6), card_width * 0.35)
        add_text_box(slide, text_left, text_top, title_w, text_h,
                     f["title"], font_size=9, font_color=DARK_NAVY, bold=True)

        desc_left = text_left + title_w + Inches(0.15)
        desc_w = card_width - title_w - 2 * inner_pad - Inches(0.15)
        add_text_box(slide, desc_left, text_top, desc_w, text_h,
                     f["desc"], font_size=8, font_color=MED_GRAY)

# ============================================================
# FINDING 7 — spans full width (aligned with Uppföljning bar)
# ============================================================
f7_y = findings_start_y + 6 * (finding_row_h + finding_gap)
f7_left = uppf_left
f7_width = uppf_width
f7_h = finding_row_h

# Card
card7 = add_rounded_rect(slide, f7_left, f7_y, f7_width, f7_h,
                          LIGHT_GRAY, border_color=DARK_NAVY, border_width=Pt(1.2))
card7.adjustments[0] = 0.06

# Number circle
num_y7 = f7_y + (f7_h - num_circle_size) / 2
circle7 = slide.shapes.add_shape(MSO_SHAPE.OVAL, NUM_MARGIN_LEFT, num_y7, num_circle_size, num_circle_size)
circle7.fill.solid()
circle7.fill.fore_color.rgb = DARK_NAVY
circle7.line.fill.background()
tf = circle7.text_frame
tf.paragraphs[0].text = finding_7["num"]
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = "Calibri"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
bodyPr = tf._txBody.find(qn('a:bodyPr'))
bodyPr.set('anchor', 'ctr')
bodyPr.set('lIns', '0')
bodyPr.set('rIns', '0')
bodyPr.set('tIns', '0')
bodyPr.set('bIns', '0')

# Text for finding 7 — wide card, title left + description right
inner_pad = Inches(0.15)
text_left7 = f7_left + inner_pad
text_top7 = f7_y + Inches(0.06)
text_h7 = f7_h - Inches(0.12)
title_w7 = Inches(3.6)

add_text_box(slide, text_left7, text_top7, title_w7, text_h7,
             finding_7["title"], font_size=9, font_color=DARK_NAVY, bold=True)

desc_left7 = text_left7 + title_w7 + Inches(0.15)
desc_w7 = f7_width - title_w7 - 2 * inner_pad - Inches(0.15)
add_text_box(slide, desc_left7, text_top7, desc_w7, text_h7,
             finding_7["desc"], font_size=8, font_color=MED_GRAY)


# ============================================================
# FOOTER
# ============================================================
footer_y = Inches(7.1)
sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), footer_y, Inches(11.93), Pt(0.75))
sep2.fill.solid()
sep2.fill.fore_color.rgb = BORDER_LINE
sep2.line.fill.background()

add_text_box(slide, Inches(0.7), footer_y + Inches(0.05), Inches(4), Inches(0.25),
             "Konfidentiellt  |  Processkartläggning P2P", font_size=7, font_color=MED_GRAY)

add_text_box(slide, Inches(10.5), footer_y + Inches(0.05), Inches(2.2), Inches(0.25),
             "Källa: Intern analys", font_size=7, font_color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


# ============================================================
# SAVE
# ============================================================
output_path = "/home/user/coachcraft-app/P2P_Processkarta_Findings.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
